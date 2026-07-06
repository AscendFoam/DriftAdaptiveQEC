from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
import zipfile

ROOT = Path(r"D:\Codes\Quantum\DriftAdaptiveQEC")
OUT = ROOT / "docs" / "paper_notes" / "ppt_output"
PPTX = OUT / "CNN_FPGA_GKP_submission_draft_two_slide_cn_v2_no_font_garble.pptx"
S1 = OUT / "slide1_v2_no_font_garble.png"
S2 = OUT / "slide2_v2_no_font_garble.png"
QA = OUT / "qa_report_v2.md"
MANIFEST = OUT / "asset_manifest_v2.md"
FIGDIR = ROOT / "docs" / "figure_assets" / "submission_draft_python_figures" / "outputs"
FIG1 = FIGDIR / "fig01_dual_loop_architecture.png"
FIG2 = FIGDIR / "fig02_main_software_hil_results.png"
FONT = r"C:\Windows\Fonts\msyh.ttc"
FONT_BOLD = r"C:\Windows\Fonts\msyhbd.ttc"
W, H = 1920, 1080

COL = {
    "bg": (248, 249, 247),
    "ink": (28, 34, 38),
    "muted": (88, 98, 105),
    "line": (205, 214, 219),
    "white": (255, 255, 255),
    "green": (41, 121, 91),
    "blue": (48, 94, 145),
    "teal": (29, 112, 121),
    "amber": (176, 113, 33),
    "red": (155, 75, 67),
    "purple": (105, 86, 156),
    "lg": (225, 239, 232),
    "lb": (226, 235, 246),
    "la": (246, 237, 222),
    "lr": (244, 229, 226),
    "lp": (238, 235, 246),
    "panel": (255, 255, 255),
    "strip": (235, 240, 238),
}


def font(size, bold=False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT, size)


def text_size(draw, text, fnt):
    box = draw.textbbox((0, 0), text, font=fnt)
    return box[2] - box[0], box[3] - box[1]


def wrap(draw, text, fnt, max_w):
    lines, cur = [], ""
    for ch in text:
        if ch == "\n":
            lines.append(cur)
            cur = ""
            continue
        candidate = cur + ch
        if text_size(draw, candidate, fnt)[0] <= max_w or not cur:
            cur = candidate
        else:
            lines.append(cur)
            cur = ch
    if cur:
        lines.append(cur)
    return lines


def draw_text(draw, xy, text, size=28, color=None, bold=False, max_w=None, line_gap=6, center=False):
    color = color or COL["ink"]
    fnt = font(size, bold)
    x, y = xy
    lines = text.split("\n") if max_w is None else wrap(draw, text, fnt, max_w)
    for line in lines:
        w, h = text_size(draw, line, fnt)
        xx = x + (max_w - w) // 2 if center and max_w else x
        draw.text((xx, y), line, font=fnt, fill=color)
        y += h + line_gap
    return y


def round_rect(draw, box, fill, outline=None, width=2, radius=16):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def paste_fit(canvas, path, box):
    img = Image.open(path).convert("RGB")
    x1, y1, x2, y2 = box
    bw, bh = x2 - x1, y2 - y1
    scale = min(bw / img.width, bh / img.height)
    nw, nh = int(img.width * scale), int(img.height * scale)
    img = img.resize((nw, nh), Image.LANCZOS)
    canvas.paste(img, (x1 + (bw - nw) // 2, y1 + (bh - nh) // 2))


def arrow(draw, p1, p2, color, width=5):
    draw.line([p1, p2], fill=color, width=width)
    x1, y1 = p1
    x2, y2 = p2
    if x2 >= x1:
        pts = [(x2, y2), (x2 - 18, y2 - 10), (x2 - 18, y2 + 10)]
    else:
        pts = [(x2, y2), (x2 + 18, y2 - 10), (x2 + 18, y2 + 10)]
    draw.polygon(pts, fill=color)


def metric_panel(draw, box, title, main, note, fill, accent):
    x1, y1, x2, y2 = box
    round_rect(draw, box, fill=fill, outline=accent, width=2, radius=12)
    draw.rectangle((x1, y1, x1 + 10, y2), fill=accent)
    draw_text(draw, (x1 + 28, y1 + 18), title, 22, accent, True, max_w=x2 - x1 - 44)
    draw_text(draw, (x1 + 28, y1 + 55), main, 30, COL["ink"], True, max_w=x2 - x1 - 44)
    draw_text(draw, (x1 + 28, y1 + 105), note, 20, COL["muted"], False, max_w=x2 - x1 - 44, line_gap=5)


def stage(draw, box, num, title, body, fill, accent):
    x1, y1, x2, y2 = box
    round_rect(draw, box, fill=fill, outline=accent, width=2, radius=18)
    draw.ellipse((x1 + 18, y1 + 18, x1 + 58, y1 + 58), fill=accent)
    draw_text(draw, (x1 + 31, y1 + 21), str(num), 19, COL["white"], True)
    draw_text(draw, (x1 + 72, y1 + 20), title, 24, COL["ink"], True, max_w=x2 - x1 - 90)
    draw_text(draw, (x1 + 22, y1 + 72), body, 18, COL["muted"], False, max_w=x2 - x1 - 44, line_gap=4)


def build_slide_1():
    im = Image.new("RGB", (W, H), COL["bg"])
    d = ImageDraw.Draw(im)
    draw_text(d, (70, 45), "项目优势：把漂移自适应压缩成可验证的快路径", 42, COL["ink"], True)
    draw_text(d, (72, 105), "基于投稿稿草案；当前证据边界是 software-HIL 和软件验证，不是真板结果", 22, COL["muted"])
    round_rect(d, (70, 175, 980, 750), COL["white"], COL["line"], 2, 4)
    paste_fit(im, FIG2, (100, 200, 950, 682))
    draw_text(d, (110, 697), "来源：投稿稿 Fig. 2。n=2 描述性 software-HIL 结果；不是置信区间或显著性检验。", 18, COL["muted"], max_w=820)
    metric_panel(d, (1030, 172, 1815, 315), "解码延迟 / 效率", "实时面只做 K*s+b", "每 shot 为 2x2 仿射计算和裁剪；4 次乘法、4 次加法、6 个状态量。", COL["lg"], COL["green"])
    metric_panel(d, (1030, 342, 1815, 485), "LER", "四个漂移场景均为最低均值", "相对 UKF 的描述性均值下降约 1.8% 到 2.9%，平均约 2.3%。", COL["lb"], COL["blue"])
    metric_panel(d, (1030, 512, 1815, 655), "保真度", "只到边界代理指标", "稿件使用 residual-boundary surrogate；不能写成有限能量逻辑通道保真度。", COL["la"], COL["amber"])
    metric_panel(d, (1030, 682, 1815, 825), "噪声适应能力", "慢环更新 K 和 b", "利用最近 syndrome 统计跟踪静态、线性、阶跃和周期漂移；仍不是 holdout 鲁棒性证明。", COL["lr"], COL["red"])
    d.rectangle((70, 895, 1815, 985), fill=COL["strip"])
    draw_text(d, (110, 915), "汇报口径：本项目优势不是“CNN 全面胜出”，而是把实时纠错固定为低成本、可量化、可回放的仿射快路径。", 28, COL["ink"], True, max_w=1660, center=True)
    draw_text(d, (990, 1012), "边界：没有真板执行成功、没有实测 FPGA latency/resource/power、没有 deployment closure。", 18, COL["muted"], max_w=820)
    im.save(S1, quality=96)


def build_slide_2():
    im = Image.new("RGB", (W, H), COL["bg"])
    d = ImageDraw.Draw(im)
    draw_text(d, (70, 45), "工程闭环：FPGA 快路径实时纠错，慢环只提交新参数", 42, COL["ink"], True)
    draw_text(d, (72, 105), "核心思想：学习和统计估计不进入 per-shot critical path；实时面只消费已经提交的参数。", 22, COL["muted"])
    round_rect(d, (70, 170, 700, 455), COL["white"], COL["line"], 2, 4)
    paste_fit(im, FIG1, (95, 190, 675, 410))
    draw_text(d, (105, 417), "来源：投稿稿 Fig. 1，dual-loop architecture。", 18, COL["muted"], max_w=560)
    round_rect(d, (760, 170, 1230, 455), COL["lg"], COL["green"], 2, 16)
    draw_text(d, (790, 205), "当前已经支撑", 28, COL["green"], True)
    draw_text(d, (790, 255), "software-HIL 漂移场景\n仿射 fast-path 成本模型\nQ4.20 软件定点一致性\nruntime counters 和 source data", 22, COL["ink"], max_w=395, line_gap=8)
    draw_text(d, (790, 405), "不等于真板 HIL 或部署闭环", 19, COL["muted"])
    round_rect(d, (1290, 170, 1765, 455), COL["lr"], COL["red"], 2, 16)
    draw_text(d, (1320, 205), "后续上板必须补", 28, COL["red"], True)
    draw_text(d, (1320, 255), "device path 和 board logs\nbitstream 或 RTL hash\nDMA / MMIO 证据\nlatency / resource / power", 22, COL["ink"], max_w=390, line_gap=8)
    draw_text(d, (1320, 405), "当前稿件只定义测量要求", 19, COL["muted"])
    boxes = []
    x0, y0, bw, bh, gap = 80, 585, 270, 145, 36
    data = [
        ("syndrome 输入", "GKP 读出连续 syndrome，保留模拟量信息。", COL["lb"], COL["blue"]),
        ("参数 bank", "读取已提交的 K 和 b，旧 bank 可回退。", COL["lg"], COL["green"]),
        ("FPGA 快路径", "执行 Delta = K*s + b，随后裁剪和限幅。", (232, 241, 242), COL["teal"]),
        ("纠错输出", "施加位移纠错，并记录 counters 与事件。", COL["la"], COL["amber"]),
        ("慢环估计", "CNN、teacher 或 statcalib 从窗口统计估计新参数。", COL["lp"], COL["purple"]),
        ("提交控制", "stage、validate、commit；异常则 hold 或 rollback。", COL["lr"], COL["red"]),
    ]
    for i, item in enumerate(data):
        x = x0 + i * (bw + gap)
        boxes.append((x, y0, x + bw, y0 + bh))
        stage(d, boxes[-1], i + 1, *item)
        if i > 0:
            arrow(d, (boxes[i - 1][2] + 8, y0 + 72), (x - 8, y0 + 72), COL["muted"], 4)
    d.line([(boxes[-1][2] - 20, 755), (boxes[-1][2] - 20, 845), (boxes[1][0] + 20, 845), (boxes[1][0] + 20, 755)], fill=COL["red"], width=4)
    d.polygon([(boxes[1][0] + 20, 755), (boxes[1][0] + 8, 777), (boxes[1][0] + 32, 777)], fill=COL["red"])
    draw_text(d, (620, 862), "慢环只更新参数；实时纠错仍保持同一个简单快路径。", 24, COL["red"], True, max_w=700, center=True)
    draw_text(d, (930, 1012), "边界：board_backend 仍是 placeholder；本页不是板级执行成功声明。", 18, COL["muted"], max_w=820)
    im.save(S2, quality=96)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for img_path in (S1, S2):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX)
    prs2 = Presentation(PPTX)
    issues = []
    for si, slide in enumerate(prs2.slides, 1):
        for sh in slide.shapes:
            if sh.left < 0 or sh.top < 0 or sh.left + sh.width > prs2.slide_width or sh.top + sh.height > prs2.slide_height:
                issues.append(f"slide {si}: shape out of bounds")
    with zipfile.ZipFile(PPTX, "r") as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
    QA.write_text(
        f"""# QA Report v2

- PPTX: `{PPTX.relative_to(ROOT)}`
- Slide images: `{S1.relative_to(ROOT)}`, `{S2.relative_to(ROOT)}`
- Slide count: {len(prs2.slides)}
- Embedded media files: {len(media)}
- Main visible text is rasterized into full-slide PNGs to avoid PowerPoint CJK font substitution, mojibake, and question-mark glyph fallback.
- Special-symbol reduction: replaced subscript/Greek-heavy formula with plain `Delta = K*s + b`; avoided rare glyphs in visible text.
- Structural bounds check: {"PASS" if not issues else "; ".join(issues)}
- Evidence boundary: no real-board execution success, no measured FPGA latency/resource/power, no finite-energy logical-channel fidelity, no p-value/CI, no deployment closure.
- Rendered preview note: the delivered slide PNGs are the rendered preview source used inside the PPTX.
""",
        encoding="utf-8",
    )
    MANIFEST.write_text(
        """# Asset Manifest v2

| Asset | Source file | Slide | Use | Boundary |
| --- | --- | --- | --- | --- |
| slide1_v2_no_font_garble.png | generated from source figures and Chinese text | 1 | Full-slide raster, prevents font-garble in visible content | Summary slide only; does not upgrade evidence |
| slide2_v2_no_font_garble.png | generated from source figures and Chinese text | 2 | Full-slide raster, prevents font-garble in visible content | Engineering explanation only; no board success claim |
| Fig. 2 main software-HIL results | `docs/figure_assets/submission_draft_python_figures/outputs/fig02_main_software_hil_results.png` | 1 | LER evidence visual | n=2 descriptive software-HIL, not CI/p-value or hardware evidence |
| Fig. 1 dual-loop architecture | `docs/figure_assets/submission_draft_python_figures/outputs/fig01_dual_loop_architecture.png` | 2 | Dual-loop architecture visual | Schematic, not board measurement |

No web assets were used.
""",
        encoding="utf-8",
    )
    print(PPTX)
    print(S1)
    print(S2)
    print(QA)
    print(f"slides={len(prs2.slides)} media={len(media)} issues={len(issues)}")


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    build_slide_1()
    build_slide_2()
    build_pptx()
